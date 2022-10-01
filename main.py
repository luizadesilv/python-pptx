from pptx import presentation, Presentation
import configuration
import docx
import lorem

# --------------- CREATE PRESENTATION ----------------#
prs = Presentation()

# --------------- SLIDE 1 ----------------#
sld0 = prs.slides.add_slide(prs.slide_layouts[1])
configuration.Style.add_hp_logo(sld0)
configuration.Style.create_line(sld0)
configuration.Style.background_color(sld0)

text = 'Presentation Title'
configuration.Text.add_title(sld0, text)

# --------------- SLIDE 2 ----------------#
sld1 = prs.slides.add_slide(prs.slide_layouts[4])
configuration.Style.add_hp_logo(sld1)
configuration.Style.create_line(sld1)
configuration.Style.background_color(sld1)

text = 'Title 2 here'
configuration.Text.add_title(sld1, text)
text1 = 'Subtitle 2 here'
configuration.Text.add_text(sld1, text1, 1, 0)

output = docx.File.read_docx(r'data/text-data.docx')
topics = docx.File.split_text(output)

subtopics = topics[1].split(":")
configuration.Text.add_text(sld1, subtopics[0])
configuration.Text.add_text(sld1, subtopics[1].replace("\n\n", "\n"), 0, 2)

subtopics = topics[2].split(":")
configuration.Text.add_text(sld1, subtopics[0], 0, 3)
configuration.Text.add_text(sld1, subtopics[1].replace("\n\n", "\n"), 0, 4)

# --------------- SLIDE 3 ----------------#
sld2 = prs.slides.add_slide(prs.slide_layouts[1])
configuration.Style.add_hp_logo(sld2)
configuration.Style.create_line(sld2)
configuration.Style.background_color(sld2)

text = 'Title 3 here'
configuration.Text.add_title(sld2, text)
text1 = 'Subtitle 3 here'
configuration.Text.add_text(sld2, text1, 1, 0)
configuration.Style.add_image(sld2, 'student.jpg')

# --------------- SLIDE 4 ----------------#
sld3 = prs.slides.add_slide(prs.slide_layouts[5])
configuration.Style.add_hp_logo(sld3)
configuration.Style.create_line(sld3)
configuration.Style.background_color(sld3)

text = 'Title 4 here'
configuration.Text.add_title(sld3, text)
text1 = 'Subtitle 4 here'
configuration.Text.add_text(sld3, text1, 1, 0)

table = docx.File.data_to_table()
configuration.Style.create_chart(table, sld3)
# --------------- SLIDE 5 ----------------#
sld4 = prs.slides.add_slide(prs.slide_layouts[3])
configuration.Style.add_hp_logo(sld4)
configuration.Style.create_line(sld4)
configuration.Style.background_color(sld4)

text = 'Title 5 here'
configuration.Text.add_title(sld4, text)
text1 = 'Subtitle 5 here'
configuration.Text.add_text(sld4, text1, 1, 0)

# --------------- SAVE PPTX  ----------------#
prs.save('DemoPresentation.pptx')
