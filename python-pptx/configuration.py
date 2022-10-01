from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


class Style:

    @staticmethod
    def add_image(slide, name):
        img = f'assets/{name}'
        slide.shapes.add_picture(img, Inches(3.5), Inches(4), height=Inches(3), width=Inches(3))

    @staticmethod
    def add_hp_logo(slide):
        logo = 'assets/logo.png'
        slide.shapes.add_picture(logo, Inches(9), Inches(6.5), height=Inches(0.75), width=Inches(0.75))

    @staticmethod
    def font_style(rn, font_size=14):
        rn.font.color.rgb = RGBColor(0, 0, 0)
        rn.font.name = 'Forma DJR Display'
        rn.font.size = Pt(font_size)

    @staticmethod
    def background_color(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(102, 178, 255)

    @staticmethod
    def create_line(slide):
        shape_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(1.5), Inches(10), 2)
        shape_line.shadow.inherit = False
        line = shape_line.line
        line.color.rgb = RGBColor(0, 0, 0)

    @staticmethod
    def create_table(table, slide):
        shape = slide.shapes.add_table(7, 6, Inches(3), Inches(2), Inches(7), Inches(5))
        tab = shape.table
        for i in range(len(table)):
            for j in range(len(table)):
                cell = tab.cell(i, j)
        cell.text = 'text'

        print(cell)

    @staticmethod
    def create_chart(table, slide):
        chart_data = ChartData()
        chart_data.categories = table['Year']
        chart_data.add_series('ML', table.loc[3])
        chart_data.add_series('BigData', table.loc[4])
        chart_data.add_series('Data Science', table.loc[5])

        x, y, cx, cy = Inches(2), Inches(2), Inches(7), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart


class Text:
    @staticmethod
    def add_title(slide, text):
        slide.shapes[0].text = text
        pg = slide.shapes[0].text_frame.paragraphs[0]
        rn = pg.runs[0]
        Style.font_style(rn, 50)

    @staticmethod
    def add_text(slide, text, index_p=0, index_s=1):
        tf = slide.shapes[index_s].text_frame
        tf.add_paragraph()
        pg = slide.shapes[index_s].text_frame.paragraphs[index_p]
        rn = pg.add_run()
        rn.text = text
        Style.font_style(rn)
